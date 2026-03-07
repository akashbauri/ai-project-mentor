from PyPDF2 import PdfReader
import docx
import requests
from bs4 import BeautifulSoup
from youtube_transcript_api import YouTubeTranscriptApi
from pptx import Presentation


def extract_text(file):

    text = ""

    try:

        # -------------------------
        # PDF
        # -------------------------
        if file.type == "application/pdf":

            reader = PdfReader(file)

            for page in reader.pages:

                page_text = page.extract_text()

                if page_text:
                    text += page_text + " "


        # -------------------------
        # WORD DOCX
        # -------------------------
        elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":

            doc = docx.Document(file)

            for para in doc.paragraphs:
                text += para.text + " "


        # -------------------------
        # POWERPOINT
        # -------------------------
        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":

            presentation = Presentation(file)

            for slide in presentation.slides:

                for shape in slide.shapes:

                    if hasattr(shape, "text"):
                        text += shape.text + " "

    except Exception as e:

        text = ""

    return text



def extract_from_url(url):

    text = ""

    try:

        # -------------------------
        # YOUTUBE
        # -------------------------
        if "youtube.com" in url or "youtu.be" in url:

            if "v=" in url:
                video_id = url.split("v=")[1].split("&")[0]
            else:
                video_id = url.split("/")[-1]

            transcript = YouTubeTranscriptApi.get_transcript(video_id)

            for t in transcript:
                text += t["text"] + " "


        # -------------------------
        # WEBSITE
        # -------------------------
        else:

            response = requests.get(url, timeout=10)

            soup = BeautifulSoup(response.text, "html.parser")

            paragraphs = soup.find_all("p")

            for p in paragraphs:
                text += p.get_text() + " "

    except Exception:

        text = ""

    return text
